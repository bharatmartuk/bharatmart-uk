"""
Generate BharatMart internship report (.docx) and presentation (.pptx)
mirroring the structure of the sample Edunet/IBM report.
Fill PERSONAL DETAILS below before submitting.
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

# ═══════════════════════════════════════════════════════════
# PERSONAL DETAILS — edit these before printing / submitting
# ═══════════════════════════════════════════════════════════
INTERN_NAME = "Pavan Kumar Kunukuntla"
STUDENT_ID = "B210074"
COLLEGE = "Rajiv Gandhi University of Knowledge Technologies, Basar"
COURSE = "B. Tech. (Computer Science and Engineering)"
DEPARTMENT = "DEPARTMENT OF COMPUTER SCIENCE AND ENGINEERING"
UNIVERSITY_SHORT = "RAJIV GANDHI UNIVERSITY OF KNOWLEDGE TECHNOLOGIES"
UNIVERSITY_PLACE = "BASAR, NIRMAL (DIST)"
FACULTY_REVIEWER = "Riharika"
FACULTY_DESIGNATION = "Assistant Professor"
HOD_NAME = "Mr. B Venkat Raman"
HOD_DESIGNATION = "Assistant Professor"
INDUSTRY_MENTOR = "Uday Kumar Kadiyam"
INDUSTRY_MENTOR_TITLE = "Mentor, BharatMart.uk"
ORG_NAME = "BharatMart.uk"
ORG_WEBSITE = "https://bharatmart.uk"
GITHUB_URL = "https://github.com/bharatmartuk/bharatmart-uk"
ROLE = "Full Stack Engineer Intern"
PROJECT_TITLE = "BharatMart UK — Multi-Merchant Grocery Marketplace"
REPORT_MONTH = "JULY 2026"
START_DATE = "15th June 2026"
END_DATE = "30th July 2026"
DECLARATION_DATE = "23-07-2026"
PLACE = "Basar"

OUT_DIR = Path(__file__).resolve().parent


def set_run_font(run, size=12, bold=False, italic=False, color=None):
    run.bold = bold
    run.italic = italic
    run.font.size = Pt(size)
    run.font.name = "Times New Roman"
    r = run._element
    r.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    if color:
        run.font.color.rgb = RGBColor(*color)


def add_centered(doc, text, size=14, bold=False, space_after=6, space_before=0):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.space_before = Pt(space_before)
    run = p.add_run(text)
    set_run_font(run, size=size, bold=bold)
    return p


def add_body(doc, text, size=12, first_line_indent=True, space_after=8):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    if first_line_indent:
        p.paragraph_format.first_line_indent = Inches(0.3)
    run = p.add_run(text)
    set_run_font(run, size=size)
    return p


def add_heading_custom(doc, text, level=1):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(8)
    run = p.add_run(text)
    size = 16 if level == 1 else 13 if level == 2 else 12
    set_run_font(run, size=size, bold=True)
    return p


def add_bullet(doc, text, size=12):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.space_after = Pt(4)
    p.clear()
    run = p.add_run(text)
    set_run_font(run, size=size)
    return p


def page_break(doc):
    doc.add_page_break()


def add_simple_table(doc, headers, rows, caption=None):
    if caption:
        cap = doc.add_paragraph()
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = cap.add_run(caption)
        set_run_font(r, size=11, italic=True)
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        run = p.add_run(h)
        set_run_font(run, size=10, bold=True)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.rows[ri + 1].cells[ci]
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            set_run_font(run, size=10)
    doc.add_paragraph()


def build_report() -> Path:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1.25)
    section.right_margin = Inches(1)

    # ── Page 1: Cover ──
    add_centered(doc, "INTERNSHIP REPORT", size=22, bold=True, space_before=72)
    add_centered(doc, "on", size=12, space_before=12)
    add_centered(doc, PROJECT_TITLE, size=16, bold=True, space_before=12)
    add_centered(doc, "performed at", size=12, space_before=18)
    add_centered(doc, ORG_NAME, size=14, bold=True, space_before=6)
    add_centered(doc, f"({ORG_WEBSITE})", size=11, space_before=4)
    add_centered(doc, "Submitted by:", size=12, space_before=36)
    add_centered(doc, INTERN_NAME, size=14, bold=True, space_before=6)
    add_centered(doc, f"College : {COLLEGE}", size=11, space_before=10)
    add_centered(doc, f"Course : {COURSE}", size=11, space_before=4)
    add_centered(doc, f"ID: {STUDENT_ID}", size=11, space_before=4)
    add_centered(doc, ORG_NAME.upper(), size=14, bold=True, space_before=48)
    add_centered(doc, "i", size=10, space_before=36)

    # ── Page 2: Title page ──
    page_break(doc)
    add_centered(doc, PROJECT_TITLE.upper(), size=16, bold=True, space_before=36)
    add_centered(doc, "Full Stack Web Marketplace", size=13, bold=True, space_before=8)
    add_centered(
        doc,
        "An Internship Report submitted to",
        size=12,
        space_before=24,
    )
    add_centered(doc, COLLEGE, size=12, bold=True, space_before=6)
    add_centered(
        doc,
        "for the partial fulfillment of the requirements\nfor the award of the degree of",
        size=11,
        space_before=10,
    )
    add_centered(doc, "Bachelor of Technology", size=13, bold=True, space_before=10)
    add_centered(doc, "in", size=11, space_before=4)
    add_centered(doc, "Computer Science & Engineering", size=12, bold=True, space_before=4)
    add_centered(doc, "Submitted by", size=11, space_before=20)
    add_centered(doc, INTERN_NAME, size=13, bold=True, space_before=4)
    add_centered(doc, STUDENT_ID, size=12, space_before=2)
    add_centered(doc, "Guided and Reviewed by", size=11, space_before=18)
    add_centered(doc, "Industry Mentor", size=11, bold=True, space_before=8)
    add_centered(doc, INDUSTRY_MENTOR, size=12, bold=True, space_before=2)
    add_centered(doc, INDUSTRY_MENTOR_TITLE, size=11, space_before=2)
    add_centered(doc, "Faculty Reviewer", size=11, bold=True, space_before=12)
    add_centered(doc, FACULTY_REVIEWER, size=12, bold=True, space_before=2)
    add_centered(doc, FACULTY_DESIGNATION, size=11, space_before=2)
    add_centered(doc, "RGUKT-BASAR", size=11, space_before=2)
    add_centered(doc, DEPARTMENT, size=11, bold=True, space_before=28)
    add_centered(doc, UNIVERSITY_SHORT, size=11, bold=True, space_before=4)
    add_centered(doc, UNIVERSITY_PLACE, size=11, bold=True, space_before=4)
    add_centered(doc, REPORT_MONTH, size=12, bold=True, space_before=16)
    add_centered(doc, "ii", size=10, space_before=12)

    # ── Page 3: Certificate ──
    page_break(doc)
    add_centered(doc, DEPARTMENT, size=11, bold=True, space_before=24)
    add_centered(doc, f"{UNIVERSITY_SHORT}, BASAR", size=11, bold=True, space_before=4)
    add_centered(doc, "CERTIFICATE", size=16, bold=True, space_before=24)
    add_body(
        doc,
        f"This is to certify that the internship report entitled “{PROJECT_TITLE}” "
        f"is a bona fide record of the internship work carried out by {INTERN_NAME}, "
        f"{STUDENT_ID}, at {ORG_NAME}, during the period {START_DATE} to {END_DATE}, "
        f"in partial fulfilment of the requirements for the award of the degree of "
        f"Bachelor of Technology in Computer Science and Engineering at "
        f"Rajiv Gandhi University of Knowledge Technologies, Basar.",
        first_line_indent=False,
    )
    add_body(
        doc,
        "The work presented in this report has not been submitted elsewhere for the "
        "award of any other degree or diploma.",
        first_line_indent=False,
    )
    doc.add_paragraph()
    p = doc.add_paragraph()
    run = p.add_run("INTERNSHIP REVIEWER")
    set_run_font(run, bold=True)
    p2 = doc.add_paragraph()
    run = p2.add_run(f"{FACULTY_REVIEWER}\n{FACULTY_DESIGNATION}")
    set_run_font(run, size=11)
    p3 = doc.add_paragraph()
    run = p3.add_run("\nHEAD OF DEPARTMENT")
    set_run_font(run, bold=True)
    p4 = doc.add_paragraph()
    run = p4.add_run(f"{HOD_NAME}\n{HOD_DESIGNATION}")
    set_run_font(run, size=11)
    add_centered(doc, "iii", size=10, space_before=24)

    # ── Page 4: Company certificate placeholder ──
    page_break(doc)
    add_centered(doc, "COMPLETION CERTIFICATE", size=16, bold=True, space_before=72)
    add_body(
        doc,
        f"[Insert the official {ORG_NAME} internship completion certificate here — "
        f"scan or PDF export. Role: {ROLE}. Duration: {START_DATE} to {END_DATE}.]",
        first_line_indent=False,
    )
    add_centered(doc, "iv", size=10, space_before=200)

    # ── Page 5: Declaration ──
    page_break(doc)
    add_centered(doc, DEPARTMENT, size=11, bold=True, space_before=24)
    add_centered(doc, f"{UNIVERSITY_SHORT}, BASAR", size=11, bold=True, space_before=4)
    add_centered(doc, "DECLARATION", size=16, bold=True, space_before=24)
    add_body(
        doc,
        f"I, {INTERN_NAME}, hereby declare that the internship report entitled "
        f"“{PROJECT_TITLE}” submitted to Rajiv Gandhi University of Knowledge "
        f"Technologies, Basar in partial fulfilment of the requirements for the award "
        f"of the degree of Bachelor of Technology, is a record of original work carried "
        f"out by me during my internship at {ORG_NAME} under the guidance of "
        f"{INDUSTRY_MENTOR}. The content presented in this report represents my "
        f"independent work and has not been submitted elsewhere for the award of any "
        f"other degree or diploma.",
        first_line_indent=False,
    )
    add_body(doc, f"Place : {PLACE}", first_line_indent=False)
    add_body(doc, f"Date : {DECLARATION_DATE}", first_line_indent=False)
    add_body(doc, f"Name of the Student — ID No\n{INTERN_NAME} — {STUDENT_ID}", first_line_indent=False)
    add_centered(doc, "v", size=10, space_before=24)

    # ── Page 6: Acknowledgement ──
    page_break(doc)
    add_centered(doc, "ACKNOWLEDGEMENT", size=16, bold=True, space_before=24)
    add_body(
        doc,
        f"I would like to express my sincere gratitude to {ORG_NAME} for providing me "
        f"the opportunity to undertake my internship and to work on the customer-facing "
        f"web marketplace for BharatMart UK.",
        first_line_indent=False,
    )
    add_body(
        doc,
        f"I am deeply thankful to {INDUSTRY_MENTOR} for continuous guidance, technical "
        f"mentorship, and constructive feedback throughout the internship period.",
        first_line_indent=False,
    )
    add_body(
        doc,
        f"I also extend my heartfelt thanks to {HOD_NAME} and {FACULTY_REVIEWER} of the "
        f"Department of Computer Science and Engineering, {COLLEGE} for their academic "
        f"support and valuable suggestions during the preparation of this report.",
        first_line_indent=False,
    )
    add_body(
        doc,
        "Finally, I thank my family, friends, and colleagues for their constant "
        "encouragement and support throughout this internship.",
        first_line_indent=False,
    )
    add_body(doc, f"{INTERN_NAME} — {STUDENT_ID}", first_line_indent=False)
    add_centered(doc, "vi", size=10, space_before=24)

    # ── Abstract ──
    page_break(doc)
    add_centered(doc, "ABSTRACT", size=16, bold=True, space_before=24)
    add_body(
        doc,
        f"Building a trustworthy multi-merchant grocery marketplace for the UK Indian "
        f"diaspora requires more than a simple product catalogue. Customers expect "
        f"local delivery awareness, secure checkout, clear seller identity, and a "
        f"polished mobile-friendly experience. This report presents the design, "
        f"architecture, and implementation of {PROJECT_TITLE}, developed during a "
        f"full-stack internship at {ORG_NAME} from {START_DATE} to {END_DATE}.",
        first_line_indent=False,
    )
    add_body(
        doc,
        "The system is implemented as a pnpm/Turborepo monorepo with three Next.js 15 "
        "applications (customer web, merchant portal, and admin console) sharing Prisma "
        "PostgreSQL models, Auth.js authentication, Zod validation, and a domain "
        "services layer. The internship focus was ownership of the customer storefront: "
        "catalog and search with fuzzy relevance ranking, cart and wishlist, UK postcode "
        "delivery-area filtering, multi-merchant checkout with Stripe PaymentIntents "
        "and cash-on-delivery, Cloudinary-backed media, and production deployment on "
        "Vercel.",
        first_line_indent=False,
    )
    add_body(
        doc,
        "Beyond feature delivery, the work emphasises production engineering: role-based "
        "auth across apps, rate limiting, seed/demo data pipelines, admin CMS for "
        "homepage banners and categories, merchant verification flows, and continuous "
        "storefront polish (mobile navigation, sticky filters, branded UI). This report "
        "documents problem analysis, architecture, implementation timeline, evaluation "
        "of shipped outcomes, and future enhancements.",
        first_line_indent=False,
    )
    add_centered(doc, "vii", size=10, space_before=24)

    # ── Contents ──
    page_break(doc)
    add_centered(doc, "Contents", size=16, bold=True, space_before=12)
    contents = [
        ("Certificate", "iii"),
        ("Company Certificate", "iv"),
        ("Declaration", "v"),
        ("Acknowledgement", "vi"),
        ("Abstract", "vii"),
        ("List of Figures", "x"),
        ("List of Tables", "xi"),
        ("List of Abbreviations", "xii"),
        ("1 INTRODUCTION", "1"),
        ("1.1 About the Organisation", "1"),
        ("1.2 Introduction to the Project", "1"),
        ("1.3 Problem Statement", "1"),
        ("1.4 Objectives of the Project", "2"),
        ("1.5 Scope of the Project", "2"),
        ("1.6 Organisation of the Report", "2"),
        ("2 INTERNSHIP LEARNINGS AND WEEKLY WORK", "3"),
        ("2.1 Full-Stack and Platform Skills", "3"),
        ("2.2 UI/UX and Product Engineering", "3"),
        ("2.3 Integrations and Deployment", "4"),
        ("2.4 Weekly Timeline (15 Jun – 23 Jul 2026)", "4"),
        ("3 PROJECT ARCHITECTURE AND SYSTEM DESIGN", "5"),
        ("3.1 System Architecture Overview", "5"),
        ("3.2 Technology Stack", "5"),
        ("3.3 Domain Model and Data Layer", "6"),
        ("3.4 Customer Experience Design", "7"),
        ("4 IMPLEMENTATION AND DEPLOYMENT", "8"),
        ("4.1 Storefront Features Implemented", "8"),
        ("4.2 Checkout, Payments, and Orders", "9"),
        ("4.3 Search, Media, and Admin Support", "10"),
        ("4.4 Deployment on Vercel", "11"),
        ("5 RESULTS AND EVALUATION", "12"),
        ("5.1 Deliverables Achieved", "12"),
        ("5.2 Qualitative Evaluation", "13"),
        ("6 CONCLUSION AND FUTURE ENHANCEMENTS", "14"),
        ("6.1 Conclusion", "14"),
        ("6.2 Future Enhancements", "14"),
        ("References", "15"),
    ]
    for title, page in contents:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        run = p.add_run(f"{title}\t{page}")
        set_run_font(run, size=11, bold=title[0].isdigit() and ". " not in title[:4] or title.startswith("References"))
    add_centered(doc, "viii", size=10, space_before=12)

    # ── List of Figures / Tables / Abbreviations ──
    page_break(doc)
    add_centered(doc, "List of Figures", size=14, bold=True)
    for line in [
        "3.1 High-level monorepo architecture (Web / Merchant / Admin)",
        "3.2 Customer journey: Browse → Location → Cart → Checkout",
        "4.1 Storefront homepage composition (hero, categories, merchants)",
        "4.2 Product detail actions: Add to cart and Add to favourites",
        "4.3 UK postcode gate and delivery filtering flow",
    ]:
        add_body(doc, line, first_line_indent=False, space_after=4)
    add_centered(doc, "x", size=10, space_before=12)

    page_break(doc)
    add_centered(doc, "List of Tables", size=14, bold=True)
    for line in [
        "2.1 Weekly internship timeline and outcomes",
        "3.1 Technology stack and roles",
        "5.1 Mapping of objectives to delivered features",
    ]:
        add_body(doc, line, first_line_indent=False, space_after=4)
    add_centered(doc, "xi", size=10, space_before=12)

    page_break(doc)
    add_centered(doc, "LIST OF ABBREVIATIONS", size=14, bold=True, space_before=12)
    abbr = [
        ("API", "Application Programming Interface"),
        ("Auth.js", "Authentication library (Auth.js / NextAuth v5)"),
        ("CI/CD", "Continuous Integration / Continuous Deployment"),
        ("CMS", "Content Management System"),
        ("COD", "Cash on Delivery"),
        ("CSS", "Cascading Style Sheets"),
        ("GBP", "Great British Pound"),
        ("JWT", "JSON Web Token"),
        ("ORM", "Object-Relational Mapping"),
        ("PDF", "Portable Document Format"),
        ("REST", "Representational State Transfer"),
        ("SEO", "Search Engine Optimisation"),
        ("SSR", "Server-Side Rendering"),
        ("UI/UX", "User Interface / User Experience"),
        ("UK", "United Kingdom"),
    ]
    add_simple_table(doc, ["Abbreviation", "Full Form"], abbr)
    add_centered(doc, "xii", size=10, space_before=8)

    # ═══ Chapter 1 ═══
    page_break(doc)
    add_heading_custom(doc, "Chapter 1", 1)
    add_heading_custom(doc, "INTRODUCTION", 1)
    add_heading_custom(doc, "1.1 About the Organisation", 2)
    add_body(
        doc,
        f"{ORG_NAME} is a UK-focused online marketplace connecting shoppers with "
        f"independent Indian grocery and homemade-food merchants. The platform aims "
        f"to make authentic regional products discoverable with reliable local delivery "
        f"expectations, transparent seller storefronts, and a modern commerce experience. "
        f"During the internship, work was carried out on the Product team with primary "
        f"ownership of the customer web application, supported by shared packages used "
        f"by merchant and admin portals.",
    )
    add_heading_custom(doc, "1.2 Introduction to the Project", 2)
    add_body(
        doc,
        f"The internship project is the end-to-end development and hardening of the "
        f"BharatMart customer marketplace. The product is organised as a monorepo "
        f"containing apps/web (storefront), apps/merchant (seller portal), and "
        f"apps/admin (operations console), plus shared packages for database, services, "
        f"auth, UI, validation, and utilities. The public codebase is available at "
        f"{GITHUB_URL}.",
    )
    add_heading_custom(doc, "1.3 Problem Statement", 2)
    add_body(
        doc,
        "Customers shopping for specialty Indian groceries in the UK face fragmented "
        "discovery across WhatsApp groups and single-shop websites. Merchants need "
        "onboarding, catalogue tools, and order handling; operators need verification "
        "and homepage CMS. A single marketplace must therefore solve: (i) product "
        "discovery with relevant search, (ii) delivery eligibility by postcode, "
        "(iii) multi-merchant cart and checkout, (iv) trusted payments, and "
        "(v) consistent branding and mobile UX.",
    )
    add_heading_custom(doc, "1.4 Objectives of the Project", 2)
    for obj in [
        "Design and implement the customer storefront UI/UX aligned with BharatMart branding.",
        "Implement catalogue browsing, fuzzy search ranking, product detail, cart, and wishlist.",
        "Introduce UK postcode-based delivery-area filtering with a guest soft-gate.",
        "Complete checkout with address selection, Stripe and COD payment paths, and multi-merchant order creation.",
        "Integrate Cloudinary media, Auth.js sessions, and Vercel production deployments.",
        "Support marketplace quality via admin CMS hooks, seed data, and production bug fixes.",
    ]:
        add_bullet(doc, obj)

    add_heading_custom(doc, "1.5 Scope of the Project", 2)
    add_body(
        doc,
        "In scope: customer web features, shared services/repositories used by the "
        "storefront, related admin/merchant flows required for a working marketplace "
        "demo, and deployment fixes. Out of scope for this report’s primary claim: "
        "trained machine-learning models, LLM chatbots, and agentic AI product features "
        "(not present in the shipped codebase). Search relevance uses heuristic fuzzy "
        "ranking rather than learned embeddings.",
    )
    add_heading_custom(doc, "1.6 Organisation of the Report", 2)
    add_body(
        doc,
        "Chapter 2 summarises internship learnings and the weekly timeline through "
        "23 July 2026. Chapter 3 describes architecture and design. Chapter 4 details "
        "implementation and deployment. Chapter 5 evaluates outcomes. Chapter 6 "
        "concludes and lists future enhancements.",
    )

    # ═══ Chapter 2 ═══
    page_break(doc)
    add_heading_custom(doc, "Chapter 2", 1)
    add_heading_custom(doc, "INTERNSHIP LEARNINGS AND WEEKLY WORK", 1)
    add_heading_custom(doc, "2.1 Full-Stack and Platform Skills", 2)
    add_body(
        doc,
        "The internship strengthened practical full-stack engineering on a modern "
        "TypeScript monorepo: Next.js App Router and Server Actions, Prisma with "
        "PostgreSQL, Auth.js credentials/OAuth, Zod validation, Zustand client state "
        "for cart/wishlist, and a server-only services layer separating repositories "
        "from UI.",
    )
    add_heading_custom(doc, "2.2 UI/UX and Product Engineering", 2)
    add_body(
        doc,
        "Significant effort went into product UX: branded header and category navigation, "
        "hero carousel, mobile bottom navigation, sticky product filters, postcode gate "
        "and location chip, auth-gated cart/favourites with resume-after-login, and "
        "replacing “chat with seller” on product pages with Add to favourites.",
    )
    add_heading_custom(doc, "2.3 Integrations and Deployment", 2)
    add_body(
        doc,
        "Hands-on integrations included Stripe PaymentIntents and webhooks, Cloudinary "
        "uploads for product and merchant documents, Resend transactional email patterns, "
        "optional Upstash rate limiting, and multi-project Vercel deployment with Prisma "
        "engine bundling and Auth URL cookie-prefix fixes.",
    )
    add_heading_custom(doc, "2.4 Weekly Timeline (15 Jun – 23 Jul 2026)", 2)
    add_simple_table(
        doc,
        ["Period", "Focus", "Key outcomes"],
        [
            [
                "15–30 Jun",
                "Foundation",
                "Requirements, stack choice, UX direction, monorepo scaffolding",
            ],
            [
                "1–17 Jul",
                "Core build",
                "Auth, catalog, cart/checkout, shared packages, portals shells",
            ],
            [
                "18 Jul",
                "Ship",
                "Monorepo milestone; Auth/Prisma/Vercel production stabilisation",
            ],
            [
                "19 Jul",
                "CMS & demo",
                "Homepage carousel admin, demo reseed, homepage cleanup",
            ],
            [
                "20 Jul",
                "UX & discovery",
                "Mobile nav, wishlist, sticky filters, smarter search, branding",
            ],
            [
                "21 Jul",
                "Content & media",
                "Info pages, categories expansion, Cloudinary product images",
            ],
            [
                "22 Jul",
                "Ops polish",
                "Marketplace admin, bulk products, rate limits, favicon, uploads",
            ],
            [
                "23 Jul",
                "Location UX",
                "Postcode delivery gate, header polish, favourites CTA on PDP",
            ],
            [
                "24–30 Jul",
                "Close-out",
                "QA, remaining commits, demo prep, report & presentation",
            ],
        ],
        caption="Table 2.1: Weekly internship timeline and outcomes",
    )

    # ═══ Chapter 3 ═══
    page_break(doc)
    add_heading_custom(doc, "Chapter 3", 1)
    add_heading_custom(doc, "PROJECT ARCHITECTURE AND SYSTEM DESIGN", 1)
    add_heading_custom(doc, "3.1 System Architecture Overview", 2)
    add_body(
        doc,
        "BharatMart is structured as four logical layers: (1) Presentation — Next.js "
        "App Router UIs for web, merchant, and admin; (2) Application — Server Actions "
        "and route handlers calling domain services; (3) Domain — packages/services "
        "repositories for products, merchants, orders, payments, reviews, and categories; "
        "(4) Data — PostgreSQL via Prisma, plus blob media on Cloudinary. Cross-cutting "
        "concerns include Auth.js sessions, Zod schemas, and rate limiting.",
    )
    add_body(
        doc,
        "[Figure 3.1 placeholder: draw a diagram with Web / Merchant / Admin apps "
        "→ shared packages (auth, services, database, ui, utils) → PostgreSQL + Cloudinary "
        "+ Stripe + Vercel.]",
        first_line_indent=False,
    )
    add_heading_custom(doc, "3.2 Technology Stack", 2)
    add_simple_table(
        doc,
        ["Tool / Framework", "Category", "Role in the Project"],
        [
            ["Next.js 15", "Web framework", "App Router storefront, merchant, admin apps"],
            ["TypeScript", "Language", "End-to-end typed monorepo"],
            ["Prisma + PostgreSQL", "Data", "Schema, migrations, queries"],
            ["Auth.js v5", "Auth", "Credentials/Google; role-scoped sessions"],
            ["Zod + RHF", "Validation/Forms", "Input schemas and form UX"],
            ["Zustand", "Client state", "Cart and wishlist persistence"],
            ["Stripe", "Payments", "PaymentIntents + webhook finalisation"],
            ["Cloudinary", "Media", "Product, banner, logo, document uploads"],
            ["Turborepo + pnpm", "Monorepo", "Task orchestration and workspaces"],
            ["Vercel", "Hosting", "Per-app production deployments"],
            ["Tailwind + shared UI", "Design system", "Branded components and layout"],
        ],
        caption="Table 3.1: Technology stack and roles",
    )
    add_heading_custom(doc, "3.3 Domain Model and Data Layer", 2)
    add_body(
        doc,
        "Core entities include User, Merchant, Product, Category, Banner, Order, "
        "MerchantOrder, Address, Coupon, and Review. Prices are stored in GBP pence. "
        "Merchants declare deliveryPostcodes; product/merchant queries can filter by "
        "customer postcode so the catalogue reflects deliverable sellers. Checkout "
        "creates one customer Order that fans out into per-merchant MerchantOrders.",
    )
    add_heading_custom(doc, "3.4 Customer Experience Design", 2)
    add_body(
        doc,
        "The first-visit experience optionally captures a UK postcode (soft gate with "
        "skip). Logged-in users derive location from saved addresses. Search combines "
        "autocomplete with heuristic fuzzy ranking (prefix/token/subsequence scoring) "
        "and popularity-based recommendations when the query is empty. Product pages "
        "expose Add to cart and Add to favourites rather than WhatsApp seller chat.",
    )
    add_body(
        doc,
        "[Figure 3.2 placeholder: customer journey flowchart Browse → Postcode → "
        "Filter merchants/products → Cart → Auth if needed → Checkout → Stripe/COD → "
        "Order confirmation.]",
        first_line_indent=False,
    )

    # ═══ Chapter 4 ═══
    page_break(doc)
    add_heading_custom(doc, "Chapter 4", 1)
    add_heading_custom(doc, "IMPLEMENTATION AND DEPLOYMENT", 1)
    add_heading_custom(doc, "4.1 Storefront Features Implemented", 2)
    add_body(
        doc,
        "Implemented storefront capabilities include homepage hero carousel (CMS-driven "
        "banners), category grid, featured merchants, product listing with filters/sort/"
        "pagination, product detail with image gallery, related products, reviews listing, "
        "mobile navigation, wishlist page, account/address areas, and info pages "
        "(About, Contact, Privacy, Terms).",
    )
    add_heading_custom(doc, "4.2 Checkout, Payments, and Orders", 2)
    add_body(
        doc,
        "Checkout is a multi-step client flow (Address → Payment → Review). Payments "
        "support Stripe PaymentIntents and COD. On successful payment webhook or COD "
        "confirmation, order finalisation creates merchant sub-orders for fulfilment. "
        "Coupons can be validated on the cart. Auth gates protect cart/wishlist actions "
        "with pending-action resume after login.",
    )
    add_heading_custom(doc, "4.3 Search, Media, and Admin Support", 2)
    add_body(
        doc,
        "Header search calls a suggest API returning autocomplete or recommended "
        "products. Product images and merchant documents are stored on Cloudinary "
        "rather than the git repository. Admin tooling supports merchant verification "
        "review, homepage carousel editing, category marketplace management, and order "
        "inspection; merchant tooling includes product CRUD, CSV bulk import, and "
        "delivery postcode settings.",
    )
    add_heading_custom(doc, "4.4 Deployment on Vercel", 2)
    add_body(
        doc,
        "Web, merchant, and admin deploy as separate Vercel projects sharing one "
        "database. Production work included Prisma query-engine packaging in the "
        "monorepo, AUTH_SECRET/Turbo env passthrough, AUTH_URL cookie-prefix alignment "
        "to stop login bounce loops, ESLint rule registration for Next.js, and "
        "Cloudinary-backed upload APIs suitable for serverless hosts.",
    )
    add_body(
        doc,
        f"Source repository: {GITHUB_URL}",
        first_line_indent=False,
    )

    # ═══ Chapter 5 ═══
    page_break(doc)
    add_heading_custom(doc, "Chapter 5", 1)
    add_heading_custom(doc, "RESULTS AND EVALUATION", 1)
    add_heading_custom(doc, "5.1 Deliverables Achieved", 2)
    add_simple_table(
        doc,
        ["Objective", "Delivered outcome"],
        [
            ["Storefront UI/UX", "Branded responsive marketplace with mobile nav"],
            ["Discovery", "Catalog, filters, fuzzy search, recommendations"],
            ["Delivery awareness", "Postcode gate + merchant delivery filtering"],
            ["Commerce", "Cart, wishlist, multi-merchant checkout, Stripe/COD"],
            ["Media & CMS", "Cloudinary assets; admin banner/category tools"],
            ["Production", "Vercel deploys; auth/media/build fixes"],
        ],
        caption="Table 5.1: Mapping of objectives to delivered features",
    )
    add_heading_custom(doc, "5.2 Qualitative Evaluation", 2)
    add_body(
        doc,
        "Success is evaluated qualitatively against internship goals: a working "
        "end-to-end purchase path, coherent brand UX, deployable monorepo, and "
        "demonstrable admin/merchant support. Git history from 18–22 July shows "
        "intensive shipping (dozens of commits) culminating in production hardening; "
        "23 July work extends postcode UX and product favourites. Remaining week "
        "(24–30 July) is reserved for QA and documentation close-out.",
    )
    add_body(
        doc,
        "Limitations acknowledged: wishlist is client-persisted (Zustand) rather than "
        "fully synced to a server Wishlist model; customer review write APIs are not "
        "the focus of this internship slice; WhatsApp remains a support deep-link, "
        "not a messaging API bot.",
    )

    # ═══ Chapter 6 ═══
    page_break(doc)
    add_heading_custom(doc, "Chapter 6", 1)
    add_heading_custom(doc, "CONCLUSION AND FUTURE ENHANCEMENTS", 1)
    add_heading_custom(doc, "6.1 Conclusion", 2)
    add_body(
        doc,
        f"This internship successfully designed and implemented the BharatMart UK "
        f"customer marketplace as a production-oriented full-stack system. Working as "
        f"a {ROLE}, I owned UI/UX and engineering for discovery, location-aware "
        f"browsing, cart/wishlist, checkout/payments, and deployment polish within a "
        f"shared monorepo architecture.",
    )
    add_body(
        doc,
        "The experience provided end-to-end ownership of a real commerce product — "
        "from interface composition and domain modelling to third-party integrations "
        "and cloud deployment — strengthening readiness for industry software roles.",
    )
    add_heading_custom(doc, "6.2 Future Enhancements", 2)
    for item in [
        "Server-synced wishlist and followed-store social features using existing Prisma models.",
        "Customer review submission and moderation workflows.",
        "Stronger personalisation (learned ranking / embeddings) on top of current fuzzy search.",
        "Merchant coupon management UI to complement cart-side coupon validation.",
        "Automated end-to-end tests for checkout and postcode filtering in CI.",
        "Expanded analytics dashboards for merchants beyond current revenue charts.",
    ]:
        add_bullet(doc, item)

    # ═══ References ═══
    page_break(doc)
    add_heading_custom(doc, "REFERENCES", 1)
    refs = [
        f"1. BharatMart.uk marketplace source code. {GITHUB_URL}",
        f"2. BharatMart.uk product site. {ORG_WEBSITE}",
        "3. Next.js Documentation, Vercel. https://nextjs.org/docs",
        "4. Prisma ORM Documentation. https://www.prisma.io/docs",
        "5. Auth.js Documentation. https://authjs.dev",
        "6. Stripe Payment Intents. https://docs.stripe.com/payments/payment-intents",
        "7. Cloudinary Upload API. https://cloudinary.com/documentation",
        "8. Turborepo Documentation. https://turbo.build/repo/docs",
        "9. Vercel Deployment Documentation. https://vercel.com/docs",
        "10. Tailwind CSS Documentation. https://tailwindcss.com/docs",
    ]
    for r in refs:
        add_body(doc, r, first_line_indent=False, space_after=6)

    out = OUT_DIR / "BharatMart_Internship_Report.docx"
    doc.save(out)
    return out


def add_slide_title(slide, text, top=0.4):
    box = slide.shapes.add_textbox(PptInches(0.5), PptInches(top), PptInches(9), PptInches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptPt(28)
    p.font.bold = True
    p.font.color.rgb = PptRGB(0x1E, 0x1B, 0x16)
    p.font.name = "Calibri"


def add_bullets(slide, lines, top=1.4, size=18):
    box = slide.shapes.add_textbox(PptInches(0.7), PptInches(top), PptInches(8.6), PptInches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = PptPt(size)
        p.font.color.rgb = PptRGB(0x51, 0x45, 0x34)
        p.font.name = "Calibri"
        p.space_after = PptPt(8)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    # 1 Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Internship Presentation", top=1.8)
    box = slide.shapes.add_textbox(PptInches(0.5), PptInches(2.8), PptInches(9), PptInches(3))
    tf = box.text_frame
    for i, line in enumerate(
        [
            PROJECT_TITLE,
            f"Role: {ROLE}",
            f"Organisation: {ORG_NAME}",
            f"Duration: {START_DATE} – {END_DATE}",
            f"Presented by: {INTERN_NAME} ({STUDENT_ID})",
            COLLEGE,
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = PptPt(18 if i == 0 else 14)
        p.font.bold = i == 0
        p.font.color.rgb = PptRGB(0x7F, 0x57, 0x00)
        p.font.name = "Calibri"

    # 2 Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Agenda")
    add_bullets(
        slide,
        [
            "1. About BharatMart.uk",
            "2. Problem & objectives",
            "3. My role and ownership",
            "4. System architecture",
            "5. What I built (features)",
            "6. Timeline of work (15 Jun – 23 Jul)",
            "7. Tech stack & deployments",
            "8. Outcomes, learnings & next steps",
        ],
    )

    # 3 About
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "About the Organisation")
    add_bullets(
        slide,
        [
            "BharatMart.uk — UK marketplace for Indian groceries & homemade foods",
            "Connects customers with independent merchants",
            "Product focus: discovery, local delivery, trusted checkout",
            f"Website: {ORG_WEBSITE}",
            f"Code: {GITHUB_URL}",
        ],
    )

    # 4 Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Problem Statement")
    add_bullets(
        slide,
        [
            "Fragmented buying via WhatsApp / single-shop sites",
            "Need multi-merchant catalogue with clear seller identity",
            "Delivery must respect UK postcode coverage",
            "Checkout must handle multi-seller baskets + payments",
            "Operators need verification + homepage CMS",
        ],
    )

    # 5 Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Objectives")
    add_bullets(
        slide,
        [
            "Own BharatMart customer web end-to-end (UI + full stack)",
            "Ship browse → search → cart → wishlist → checkout",
            "Add UK postcode delivery filtering UX",
            "Integrate Stripe/COD, Cloudinary, Auth.js",
            "Deploy production apps on Vercel and harden auth/builds",
        ],
    )

    # 6 Role
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "My Role — Full Stack Engineer Intern")
    add_bullets(
        slide,
        [
            "Primary ownership: apps/web customer marketplace",
            "UI/UX: branding, header, mobile nav, product & checkout flows",
            "Backend: services/repositories, Server Actions, Prisma models",
            "Also supported: admin CMS hooks, merchant media/uploads, deploys",
            "Not claimed as product ML/agentic AI — focus is commerce engineering",
        ],
    )

    # 7 Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Architecture Overview")
    add_bullets(
        slide,
        [
            "Monorepo: Turborepo + pnpm workspaces",
            "Apps: Web (customers) · Merchant · Admin",
            "Shared: database · services · auth · ui · utils · validation",
            "Data: PostgreSQL (Prisma) · Media: Cloudinary",
            "Payments: Stripe · Hosting: Vercel (3 projects)",
        ],
    )

    # 8 Features 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "What I Built — Discovery & UX")
    add_bullets(
        slide,
        [
            "Homepage: hero carousel, categories, featured merchants",
            "Product listing: filters, sort, sticky desktop chrome",
            "Fuzzy search + autocomplete / recommended products",
            "Product detail: gallery, reviews, related items",
            "Favourites CTA (replaced chat-with-seller)",
            "Mobile bottom nav + search sheet",
        ],
    )

    # 9 Features 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "What I Built — Commerce & Location")
    add_bullets(
        slide,
        [
            "Cart + wishlist with auth-gated pending actions",
            "Checkout stepper: Address → Payment → Review",
            "Stripe PaymentIntents + COD; multi-merchant orders",
            "UK postcode gate, banner, location chip",
            "Filter merchants/products by delivery postcodes",
            "Account / addresses / info pages",
        ],
    )

    # 10 Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Timeline (to 23 Jul 2026)")
    add_bullets(
        slide,
        [
            "15–30 Jun: requirements, architecture, UX foundation",
            "1–17 Jul: core marketplace build (auth, catalog, checkout)",
            "18 Jul: monorepo milestone + Vercel/Auth production fixes",
            "19–21 Jul: CMS, demo data, mobile/wishlist, Cloudinary, info pages",
            "22 Jul: marketplace admin, bulk import, rate limits, polish",
            "23 Jul: postcode delivery UX + favourites on product page",
            "24–30 Jul: QA, docs, final demo (remaining)",
        ],
        size=16,
    )

    # 11 Tech
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Technology Stack")
    add_bullets(
        slide,
        [
            "Next.js 15 · TypeScript · Tailwind · shared UI kit",
            "Prisma · PostgreSQL · Auth.js · Zod · Zustand",
            "Stripe · Cloudinary · Resend · Upstash (optional)",
            "Turborepo · pnpm · Vercel",
        ],
    )

    # 12 Demo flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Demo Story (suggested)")
    add_bullets(
        slide,
        [
            "1. Open homepage → hero + categories",
            "2. Enter postcode → see deliverable merchants/products",
            "3. Search a product → open PDP → Add to favourites",
            "4. Add to cart → checkout → Stripe test / COD",
            "5. (Optional) Admin banner edit / merchant verification",
        ],
    )

    # 13 Outcomes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Outcomes & Learnings")
    add_bullets(
        slide,
        [
            "Shipped a real multi-merchant marketplace storefront",
            "Learned production auth, payments, media, and monorepo deploys",
            "Practised product UX under commerce constraints",
            "Improved debugging of Vercel/Prisma/Auth edge cases",
            "Gained ownership mindset: design → implement → ship → polish",
        ],
    )

    # 14 Future
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Future Enhancements")
    add_bullets(
        slide,
        [
            "Server-synced wishlist & followed stores",
            "Customer review write + moderation",
            "Learned search ranking on top of fuzzy heuristics",
            "Merchant coupon management UI",
            "Stronger E2E test coverage in CI",
        ],
    )

    # 15 Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Thank You", top=2.5)
    box = slide.shapes.add_textbox(PptInches(0.5), PptInches(3.5), PptInches(9), PptInches(2))
    tf = box.text_frame
    for i, line in enumerate(
        [
            "Questions welcome",
            f"{INTERN_NAME} — {STUDENT_ID}",
            f"{ORG_NAME} · {ROLE}",
            GITHUB_URL,
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = PptPt(16)
        p.font.color.rgb = PptRGB(0x51, 0x45, 0x34)
        p.font.name = "Calibri"

    out = OUT_DIR / "BharatMart_Internship_Presentation.pptx"
    prs.save(out)
    return out


def main():
    report = build_report()
    ppt = build_pptx()
    readme = OUT_DIR / "HOW_TO_FILL_REPORT.md"
    readme.write_text(
        f"""# BharatMart Internship Report & PPT

Generated files:
- `BharatMart_Internship_Report.docx` — same structure as your friend's sample PDF
- `BharatMart_Internship_Presentation.pptx` — explanation slides for the review

## Fill these placeholders before submitting

Edit `generate_internship_docs.py` (PERSONAL DETAILS section) and re-run:

```bash
python Report_Sample/generate_internship_docs.py
```

Required fields:
- INTERN_NAME
- STUDENT_ID
- FACULTY_REVIEWER / FACULTY_DESIGNATION
- HOD_NAME / HOD_DESIGNATION
- Confirm INDUSTRY_MENTOR (default: {INDUSTRY_MENTOR})
- Dates currently: {START_DATE} → {END_DATE}
- Declaration date: {DECLARATION_DATE}

## Role used in documents

**{ROLE}** (Full Stack) — matches shipped BharatMart web work.

## After generate

1. Insert company completion certificate image on the Company Certificate page.
2. Add screenshots into Word where figure placeholders are noted.
3. Export DOCX → PDF for final submission.
4. Practice PPT demo flow (slide “Demo Story”).
""",
        encoding="utf-8",
    )
    print("Wrote:", report)
    print("Wrote:", ppt)
    print("Wrote:", readme)


if __name__ == "__main__":
    main()
